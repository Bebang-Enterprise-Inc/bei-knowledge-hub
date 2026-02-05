"""PowerPoint (PPTX) extraction with vision processing."""

from typing import Dict, Any, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_pptx(filepath: str, process_images: bool = True) -> Dict[str, Any]:
    """
    Extract content from PowerPoint file.

    Args:
        filepath: Path to .pptx file
        process_images: If True, process images with vision (default: True)

    Returns:
        Dict with slides, speaker_notes, images, metadata
    """
    # Import vision module conditionally
    if process_images:
        from ..services.vision import process_image

    prs = Presentation(filepath)

    slides = []
    speaker_notes = []
    images = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []

        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)

            # Process images
            if process_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    vision_result = process_image(image_bytes)

                    images.append({
                        "slide_number": slide_num,
                        "type": vision_result["type"],
                        "quality_score": vision_result["quality_score"],
                        "extracted_text": vision_result["extracted_text"],
                        "should_index": vision_result["should_index"]
                    })
                except Exception as e:
                    # Log error but continue processing
                    print(f"Warning: Failed to process image on slide {slide_num}: {e}")

        slides.append({
            "slide_number": slide_num,
            "content": "\n".join(slide_content)
        })

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                speaker_notes.append({
                    "slide_number": slide_num,
                    "notes": notes_frame.text
                })

    return {
        "slides": slides,
        "speaker_notes": speaker_notes,
        "images": images,
        "metadata": {
            "total_slides": len(prs.slides),
            "images_processed": len(images)
        }
    }
