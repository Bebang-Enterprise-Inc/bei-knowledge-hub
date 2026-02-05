"""Tests for document extractors."""

import pytest
from unittest.mock import Mock, patch


# PPTX Extractor Tests
def test_pptx_extractor_processes_images():
    """PPTX extractor should process embedded images with vision."""
    from app.extractors.pptx import extract_pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Mock shape with image
    mock_image_shape = Mock()
    mock_image_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    mock_image_shape.image.blob = b"fake_image_bytes"
    mock_image_shape.text = ""  # Empty text for image shape

    mock_text_shape = Mock()
    mock_text_shape.has_text_frame = True
    mock_text_shape.text = "Slide title"
    mock_text_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX

    mock_slide = Mock()
    mock_slide.shapes = [mock_text_shape, mock_image_shape]
    mock_slide.has_notes_slide = False

    mock_prs = Mock()
    mock_prs.slides = [mock_slide]

    with patch("app.extractors.pptx.Presentation", return_value=mock_prs), \
         patch("app.services.vision.process_image") as mock_vision:

        mock_vision.return_value = {
            "type": "chart",
            "quality_score": 1.0,
            "extracted_text": "Sales chart Q1 2026",
            "should_index": True
        }

        result = extract_pptx("/fake/path.pptx", process_images=True)

        assert "images" in result
        assert len(result["images"]) == 1
        assert result["images"][0]["type"] == "chart"
        assert result["images"][0]["extracted_text"] == "Sales chart Q1 2026"
        mock_vision.assert_called_once()


def test_pptx_extractor_skips_images_when_disabled():
    """Should skip image processing if process_images=False."""
    from app.extractors.pptx import extract_pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    mock_image_shape = Mock()
    mock_image_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    mock_image_shape.text = ""  # Empty text

    mock_slide = Mock()
    mock_slide.shapes = [mock_image_shape]
    mock_slide.has_notes_slide = False

    mock_prs = Mock()
    mock_prs.slides = [mock_slide]

    with patch("app.extractors.pptx.Presentation", return_value=mock_prs):
        result = extract_pptx("/fake/path.pptx", process_images=False)

        assert result["images"] == []


# PDF Extractor Tests
def test_pdf_extractor_extracts_text():
    """PDF extractor should extract text from pages."""
    from app.extractors.pdf import extract_pdf

    mock_page = Mock()
    mock_page.extract_text.return_value = "Page 1 content"

    mock_reader = Mock()
    mock_reader.pages = [mock_page]

    with patch("app.extractors.pdf.PdfReader", return_value=mock_reader):
        result = extract_pdf("/fake/path.pdf", process_images=False)

        assert "pages" in result
        assert len(result["pages"]) == 1
        assert result["pages"][0]["content"] == "Page 1 content"


def test_pdf_extractor_processes_images():
    """PDF extractor should process embedded images."""
    from app.extractors.pdf import extract_pdf

    mock_image = Mock()
    mock_image.data = b"fake_image_bytes"

    mock_page = Mock()
    mock_page.extract_text.return_value = "Page with chart"
    mock_page.images = [mock_image]

    mock_reader = Mock()
    mock_reader.pages = [mock_page]

    with patch("app.extractors.pdf.PdfReader", return_value=mock_reader), \
         patch("app.extractors.pdf.process_image") as mock_vision:

        mock_vision.return_value = {
            "type": "chart",
            "quality_score": 1.0,
            "extracted_text": "Revenue chart",
            "should_index": True
        }

        result = extract_pdf("/fake/path.pdf", process_images=True)

        assert len(result["images"]) == 1
        assert result["images"][0]["type"] == "chart"


# DOCX Extractor Tests
def test_docx_extractor_extracts_text():
    """DOCX extractor should extract paragraphs."""
    from app.extractors.docx import extract_docx

    mock_para = Mock()
    mock_para.text = "Document paragraph 1"

    mock_doc = Mock()
    mock_doc.paragraphs = [mock_para]
    mock_doc.inline_shapes = []

    with patch("app.extractors.docx.Document", return_value=mock_doc):
        result = extract_docx("/fake/path.docx", process_images=False)

        assert "content" in result
        assert "Document paragraph 1" in result["content"]


def test_docx_extractor_processes_images():
    """DOCX extractor should process inline images."""
    from app.extractors.docx import extract_docx

    mock_inline_shape = Mock()
    mock_inline_shape.type = 3  # PICTURE type
    mock_inline_shape.image.blob = b"fake_bytes"

    mock_para = Mock()
    mock_para.text = "Text with image"

    mock_doc = Mock()
    mock_doc.paragraphs = [mock_para]
    mock_doc.inline_shapes = [mock_inline_shape]

    with patch("app.extractors.docx.Document", return_value=mock_doc), \
         patch("app.extractors.docx.process_image") as mock_vision:

        mock_vision.return_value = {
            "type": "diagram",
            "quality_score": 0.9,
            "extracted_text": "Process flow",
            "should_index": True
        }

        result = extract_docx("/fake/path.docx", process_images=True)

        assert len(result["images"]) == 1
        assert result["images"][0]["type"] == "diagram"
